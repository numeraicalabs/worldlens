"""
WorldLens — ETF Tracker router
Provides all ETF Tracker API endpoints under /api/etf/...
Uses the same WorldLens aiosqlite DB and auth system.
"""
from __future__ import annotations
import json
import io
import os
import smtplib
import logging
from datetime import datetime
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
from pathlib import Path
from typing import Optional, List

import aiosqlite
from fastapi import APIRouter, HTTPException, Depends
from fastapi.responses import FileResponse
from pydantic import BaseModel

from auth import get_current_user
from config import settings

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/api/etf", tags=["etf-tracker"])

REPORTS_DIR = Path(__file__).parent.parent / "reports"
REPORTS_DIR.mkdir(exist_ok=True)


# ── Pydantic models ────────────────────────────────────────────────────────────

class PortfolioCreate(BaseModel):
    name: str = "Portafoglio Principale"
    strategy: str = "custom"

class HoldingCreate(BaseModel):
    isin: str
    ticker: str
    name: str
    shares: float
    avgPrice: float

class AlertCreate(BaseModel):
    etfIsin: str
    etfTicker: str
    type: str = "below"
    threshold: float
    currentPrice: float = 0
    channels: List[str] = ["email", "push"]

class AlertUpdate(BaseModel):
    active: Optional[int] = None

class PostCreate(BaseModel):
    content: str
    portfolio_snapshot: Optional[str] = None

class SettingsUpdate(BaseModel):
    class Config:
        extra = "allow"

class ReportRequest(BaseModel):
    format: str = "pdf"
    type: str = "portfolio_summary"

class EmailRequest(BaseModel):
    format: str = "pdf"
    to: Optional[str] = None


# ── Helpers ────────────────────────────────────────────────────────────────────

async def get_db():
    return await aiosqlite.connect(settings.db_path)


def _row(cursor, row):
    """Convert aiosqlite Row to dict."""
    return {k: row[i] for i, k in enumerate([d[0] for d in cursor.description])}


# ── ONBOARDING (ETF Tracker compat) ───────────────────────────────────────────

@router.post("/onboarding")
async def save_onboarding(
    data: dict,
    user=Depends(get_current_user)
):
    async with aiosqlite.connect(settings.db_path) as db:
        await db.execute(
            "UPDATE users SET onboarding_done=1 WHERE id=?",
            (user["id"],)
        )
        await db.commit()
    return {"success": True}


# ── PORTFOLIOS ─────────────────────────────────────────────────────────────────

@router.get("/portfolios")
async def list_portfolios(user=Depends(get_current_user)):
    async with aiosqlite.connect(settings.db_path) as db:
        db.row_factory = aiosqlite.Row
        async with db.execute(
            "SELECT * FROM etf_portfolios WHERE user_id=? ORDER BY created_at",
            (user["id"],)
        ) as cur:
            portfolios = [dict(row) for row in await cur.fetchall()]

        for p in portfolios:
            async with db.execute(
                "SELECT * FROM etf_holdings WHERE portfolio_id=?",
                (p["id"],)
            ) as cur2:
                p["holdings"] = [dict(row) for row in await cur2.fetchall()]

    return portfolios


@router.post("/portfolios", status_code=201)
async def create_portfolio(
    data: PortfolioCreate,
    user=Depends(get_current_user)
):
    async with aiosqlite.connect(settings.db_path) as db:
        cur = await db.execute(
            "INSERT INTO etf_portfolios (user_id, name, strategy) VALUES (?,?,?)",
            (user["id"], data.name, data.strategy)
        )
        await db.commit()
        return {"id": cur.lastrowid, "name": data.name, "strategy": data.strategy}


@router.post("/portfolios/{pid}/holdings", status_code=201)
async def add_holding(
    pid: int,
    data: HoldingCreate,
    user=Depends(get_current_user)
):
    async with aiosqlite.connect(settings.db_path) as db:
        async with db.execute(
            "SELECT id FROM etf_portfolios WHERE id=? AND user_id=?",
            (pid, user["id"])
        ) as cur:
            if not await cur.fetchone():
                raise HTTPException(404, "Portfolio not found")

        cur2 = await db.execute(
            "INSERT INTO etf_holdings (portfolio_id, isin, ticker, name, shares, avg_price) "
            "VALUES (?,?,?,?,?,?)",
            (pid, data.isin, data.ticker, data.name, data.shares, data.avgPrice)
        )
        await db.commit()
        return {"id": cur2.lastrowid}


@router.delete("/holdings/{hid}")
async def delete_holding(hid: int, user=Depends(get_current_user)):
    async with aiosqlite.connect(settings.db_path) as db:
        await db.execute(
            "DELETE FROM etf_holdings WHERE id=? AND portfolio_id IN "
            "(SELECT id FROM etf_portfolios WHERE user_id=?)",
            (hid, user["id"])
        )
        await db.commit()
    return {"success": True}


# ── ALERTS ─────────────────────────────────────────────────────────────────────

@router.get("/alerts")
async def list_alerts(user=Depends(get_current_user)):
    async with aiosqlite.connect(settings.db_path) as db:
        db.row_factory = aiosqlite.Row
        async with db.execute(
            "SELECT * FROM etf_alerts WHERE user_id=? ORDER BY created_at DESC",
            (user["id"],)
        ) as cur:
            return [dict(row) for row in await cur.fetchall()]


@router.post("/alerts", status_code=201)
async def create_alert(data: AlertCreate, user=Depends(get_current_user)):
    async with aiosqlite.connect(settings.db_path) as db:
        cur = await db.execute(
            "INSERT INTO etf_alerts (user_id, etf_isin, etf_ticker, alert_type, "
            "threshold, current_price, channels) VALUES (?,?,?,?,?,?,?)",
            (user["id"], data.etfIsin, data.etfTicker, data.type,
             data.threshold, data.currentPrice, json.dumps(data.channels))
        )
        await db.commit()
        return {"id": cur.lastrowid}


@router.patch("/alerts/{aid}")
async def update_alert(aid: int, data: AlertUpdate, user=Depends(get_current_user)):
    async with aiosqlite.connect(settings.db_path) as db:
        if data.active is not None:
            await db.execute(
                "UPDATE etf_alerts SET active=? WHERE id=? AND user_id=?",
                (data.active, aid, user["id"])
            )
        else:
            # toggle
            await db.execute(
                "UPDATE etf_alerts SET active=NOT active WHERE id=? AND user_id=?",
                (aid, user["id"])
            )
        await db.commit()
    return {"success": True}


@router.delete("/alerts/{aid}")
async def delete_alert(aid: int, user=Depends(get_current_user)):
    async with aiosqlite.connect(settings.db_path) as db:
        await db.execute(
            "DELETE FROM etf_alerts WHERE id=? AND user_id=?",
            (aid, user["id"])
        )
        await db.commit()
    return {"success": True}


# ── SETTINGS ───────────────────────────────────────────────────────────────────

@router.get("/settings")
async def get_settings(user=Depends(get_current_user)):
    async with aiosqlite.connect(settings.db_path) as db:
        db.row_factory = aiosqlite.Row
        async with db.execute(
            "SELECT key, value FROM etf_settings WHERE user_id=?",
            (user["id"],)
        ) as cur:
            rows = await cur.fetchall()
    result = {}
    for row in rows:
        try:
            result[row["key"]] = json.loads(row["value"])
        except Exception:
            result[row["key"]] = row["value"]
    return result


@router.put("/settings")
async def save_settings(data: dict, user=Depends(get_current_user)):
    async with aiosqlite.connect(settings.db_path) as db:
        for key, value in data.items():
            v = json.dumps(value) if isinstance(value, (dict, list)) else str(value)
            await db.execute(
                "INSERT INTO etf_settings (user_id, key, value) VALUES (?,?,?) "
                "ON CONFLICT(user_id, key) DO UPDATE SET value=excluded.value",
                (user["id"], key, v)
            )
        await db.commit()
    return {"success": True}


# ── COMMUNITY ──────────────────────────────────────────────────────────────────

@router.get("/community/posts")
async def list_posts():
    async with aiosqlite.connect(settings.db_path) as db:
        db.row_factory = aiosqlite.Row
        async with db.execute(
            "SELECT * FROM etf_community_posts ORDER BY created_at DESC LIMIT 50"
        ) as cur:
            return [dict(row) for row in await cur.fetchall()]


@router.post("/community/posts", status_code=201)
async def create_post(data: PostCreate, user=Depends(get_current_user)):
    name = user.get("username") or user.get("email", "Utente")
    avatar = name[:2].upper()
    async with aiosqlite.connect(settings.db_path) as db:
        cur = await db.execute(
            "INSERT INTO etf_community_posts (user_id, user_name, avatar, content, portfolio_snapshot) "
            "VALUES (?,?,?,?,?)",
            (user["id"], name, avatar, data.content, data.portfolio_snapshot)
        )
        await db.commit()
        pid = cur.lastrowid
        async with db.execute(
            "SELECT * FROM etf_community_posts WHERE id=?", (pid,)
        ) as cur2:
            db.row_factory = aiosqlite.Row
            row = await cur2.fetchone()
            return dict(row) if row else {"id": pid}


@router.post("/community/posts/{pid}/like")
async def like_post(pid: int, user=Depends(get_current_user)):
    async with aiosqlite.connect(settings.db_path) as db:
        await db.execute(
            "UPDATE etf_community_posts SET likes=likes+1 WHERE id=?", (pid,)
        )
        await db.commit()
    return {"success": True}


# ── STATIC DATA ────────────────────────────────────────────────────────────────

@router.get("/data/etfs")
async def get_etfs():
    """Return the static ETF database."""
    return ETF_DATABASE


@router.get("/data/models")
async def get_models():
    """Return the portfolio model list."""
    return MODEL_PORTFOLIOS


# ── REPORTS ────────────────────────────────────────────────────────────────────

@router.post("/reports/generate")
async def generate_report(req: ReportRequest, user=Depends(get_current_user)):
    """Generate a PDF or PPTX portfolio report."""
    # Fetch user's portfolios
    async with aiosqlite.connect(settings.db_path) as db:
        db.row_factory = aiosqlite.Row
        async with db.execute(
            "SELECT * FROM etf_portfolios WHERE user_id=?", (user["id"],)
        ) as cur:
            portfolios = [dict(r) for r in await cur.fetchall()]
        for p in portfolios:
            async with db.execute(
                "SELECT * FROM etf_holdings WHERE portfolio_id=?", (p["id"],)
            ) as cur2:
                p["holdings"] = [dict(r) for r in await cur2.fetchall()]

    holdings = portfolios[0]["holdings"] if portfolios else []
    total_value = sum(h["shares"] * (h.get("current_price") or h["avg_price"]) for h in holdings)
    invested = sum(h["shares"] * h["avg_price"] for h in holdings)
    pl = total_value - invested
    username = user.get("username") or user.get("email", "Utente")

    report_id = None
    filepath = None

    if req.format == "pdf":
        filepath = _gen_pdf(username, holdings, total_value, invested, pl)
    elif req.format == "pptx":
        filepath = _gen_pptx(username, holdings, total_value, invested, pl)

    if filepath:
        async with aiosqlite.connect(settings.db_path) as db:
            cur = await db.execute(
                "INSERT INTO etf_reports (user_id, format, type, filepath) VALUES (?,?,?,?)",
                (user["id"], req.format, req.type, str(filepath))
            )
            await db.commit()
            report_id = cur.lastrowid

    return {"id": report_id, "format": req.format, "status": "ready" if filepath else "error"}


@router.get("/reports/download/{rid}")
async def download_report(rid: int, user=Depends(get_current_user)):
    async with aiosqlite.connect(settings.db_path) as db:
        db.row_factory = aiosqlite.Row
        async with db.execute(
            "SELECT * FROM etf_reports WHERE id=? AND user_id=?", (rid, user["id"])
        ) as cur:
            row = await cur.fetchone()
    if not row:
        raise HTTPException(404, "Report not found")
    fp = Path(row["filepath"])
    if not fp.exists():
        raise HTTPException(404, "Report file missing")
    media = "application/pdf" if row["format"] == "pdf" else "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    return FileResponse(str(fp), media_type=media, filename=fp.name)


@router.post("/reports/email")
async def email_report(req: EmailRequest, user=Depends(get_current_user)):
    """Generate and email the report."""
    to = req.to or user.get("email", "")
    if not to:
        raise HTTPException(400, "No email address")

    # Generate report first
    gen = await generate_report(ReportRequest(format=req.format), user)
    if not gen.get("id"):
        raise HTTPException(500, "Report generation failed")

    async with aiosqlite.connect(settings.db_path) as db:
        db.row_factory = aiosqlite.Row
        async with db.execute(
            "SELECT * FROM etf_reports WHERE id=?", (gen["id"],)
        ) as cur:
            row = await cur.fetchone()

    fp = Path(row["filepath"]) if row else None
    success = _send_email(to, fp, req.format)
    return {"success": success}


# ── PDF/PPTX generators ────────────────────────────────────────────────────────

def _gen_pdf(username: str, holdings: list, total: float, invested: float, pl: float) -> Optional[Path]:
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = REPORTS_DIR / f"portfolio_{ts}.pdf"
    try:
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.lib.units import cm

        doc = SimpleDocTemplate(str(path), pagesize=A4,
                                leftMargin=2*cm, rightMargin=2*cm,
                                topMargin=2*cm, bottomMargin=2*cm)
        styles = getSampleStyleSheet()
        story = []

        # Title
        title_style = ParagraphStyle("title", parent=styles["Title"],
                                     fontSize=24, spaceAfter=6, textColor=colors.HexColor("#0A2540"))
        story.append(Paragraph("ETF Tracker — Report Portafoglio", title_style))
        story.append(Paragraph(f"Generato il {datetime.now().strftime('%d/%m/%Y %H:%M')} · {username}",
                                styles["Normal"]))
        story.append(Spacer(1, 0.5*cm))

        # KPI table
        kpi_data = [
            ["Metrica", "Valore"],
            ["Valore totale", f"EUR {total:,.2f}"],
            ["Totale investito", f"EUR {invested:,.2f}"],
            ["P&L", f"EUR {pl:+,.2f}"],
            ["Rendimento", f"{(pl/invested*100 if invested else 0):+.2f}%"],
        ]
        kpi_table = Table(kpi_data, colWidths=[6*cm, 6*cm])
        kpi_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0A2540")),
            ("TEXTCOLOR",  (0, 0), (-1, 0), colors.white),
            ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE",   (0, 0), (-1, -1), 11),
            ("GRID",       (0, 0), (-1, -1), 0.5, colors.HexColor("#DCE3ED")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F7F9FC")]),
            ("LEFTPADDING",  (0, 0), (-1, -1), 10),
            ("RIGHTPADDING", (0, 0), (-1, -1), 10),
            ("TOPPADDING",   (0, 0), (-1, -1), 8),
            ("BOTTOMPADDING",(0, 0), (-1, -1), 8),
        ]))
        story.append(kpi_table)
        story.append(Spacer(1, 0.5*cm))

        # Holdings table
        if holdings:
            story.append(Paragraph("Posizioni aperte", styles["Heading2"]))
            story.append(Spacer(1, 0.2*cm))
            hold_data = [["Ticker", "ISIN", "Quote", "Prezzo medio", "Valore", "P&L"]]
            for h in holdings:
                price = h.get("current_price") or h["avg_price"]
                val = h["shares"] * price
                hpl = h["shares"] * (price - h["avg_price"])
                hold_data.append([
                    h["ticker"], h["isin"],
                    f"{h['shares']:.2f}", f"EUR {h['avg_price']:.2f}",
                    f"EUR {val:,.2f}", f"EUR {hpl:+,.2f}"
                ])
            hold_table = Table(hold_data, colWidths=[2*cm, 3.5*cm, 2*cm, 3*cm, 3*cm, 3*cm])
            hold_table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0A2540")),
                ("TEXTCOLOR",  (0, 0), (-1, 0), colors.white),
                ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE",   (0, 0), (-1, -1), 9),
                ("GRID",       (0, 0), (-1, -1), 0.3, colors.HexColor("#DCE3ED")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F7F9FC")]),
                ("LEFTPADDING",  (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING",   (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING",(0, 0), (-1, -1), 6),
            ]))
            story.append(hold_table)

        story.append(Spacer(1, cm))
        story.append(Paragraph(
            "Disclaimer: ETF Tracker è uno strumento di analisi e non costituisce consulenza finanziaria.",
            styles["Italic"]
        ))
        doc.build(story)
        return path
    except ImportError:
        logger.warning("reportlab not installed, skipping PDF generation")
        return None
    except Exception as e:
        logger.error("PDF generation error: %s", e)
        return None


def _gen_pptx(username: str, holdings: list, total: float, invested: float, pl: float) -> Optional[Path]:
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = REPORTS_DIR / f"portfolio_{ts}.pptx"
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        def add_text(slide, text, x, y, w, h, size=18, bold=False, color=(10,37,64), align=PP_ALIGN.LEFT):
            txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = txb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*color)
            return txb

        def add_rect(slide, x, y, w, h, color=(10,37,64)):
            shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(*color)
            shape.line.fill.background()
            return shape

        # Slide 1 — Title
        s1 = prs.slides.add_slide(blank)
        add_rect(s1, 0, 0, 13.33, 7.5, (6, 21, 40))
        add_text(s1, "ETF Tracker", 1, 1.5, 11, 1.2, 48, True, (255,255,255), PP_ALIGN.CENTER)
        add_text(s1, "Report Portafoglio", 1, 2.8, 11, 0.8, 28, False, (138,148,166), PP_ALIGN.CENTER)
        add_text(s1, f"{username}  ·  {datetime.now().strftime('%d/%m/%Y')}", 1, 3.8, 11, 0.6, 16, False, (138,148,166), PP_ALIGN.CENTER)

        # Slide 2 — KPIs
        s2 = prs.slides.add_slide(blank)
        add_rect(s2, 0, 0, 13.33, 1.2, (10, 37, 64))
        add_text(s2, "Riepilogo Portafoglio", 0.3, 0.2, 12, 0.8, 24, True, (255,255,255))
        kpis = [
            ("Valore totale", f"EUR {total:,.2f}", (15,123,63)),
            ("Totale investito", f"EUR {invested:,.2f}", (10,37,64)),
            ("P&L", f"EUR {pl:+,.2f}", (15,123,63) if pl >= 0 else (180,35,24)),
            ("Rendimento", f"{(pl/invested*100 if invested else 0):+.2f}%", (15,123,63) if pl >= 0 else (180,35,24)),
        ]
        for i, (lbl, val, col) in enumerate(kpis):
            x = 0.5 + (i % 2) * 6.4
            y = 1.8 + (i // 2) * 2.2
            add_rect(s2, x, y, 6, 1.8, (247,249,252))
            add_text(s2, lbl, x+0.2, y+0.2, 5.6, 0.5, 12, False, (138,148,166))
            add_text(s2, val, x+0.2, y+0.7, 5.6, 0.8, 22, True, col)

        # Slide 3 — Holdings
        if holdings:
            s3 = prs.slides.add_slide(blank)
            add_rect(s3, 0, 0, 13.33, 1.2, (10,37,64))
            add_text(s3, "Posizioni aperte", 0.3, 0.2, 12, 0.8, 24, True, (255,255,255))
            headers = ["Ticker", "Quote", "Prezzo medio", "Valore attuale", "P&L"]
            col_w = [1.8, 1.8, 2.4, 2.4, 2.4]
            x_starts = [0.3]
            for w in col_w[:-1]: x_starts.append(x_starts[-1] + w)
            for i, h_txt in enumerate(headers):
                add_text(s3, h_txt, x_starts[i], 1.4, col_w[i], 0.4, 11, True, (10,37,64))
            for ri, h in enumerate(holdings[:8]):
                price = h.get("current_price") or h["avg_price"]
                val = h["shares"] * price
                hpl = h["shares"] * (price - h["avg_price"])
                row_vals = [h["ticker"], f"{h['shares']:.2f}", f"EUR {h['avg_price']:.2f}", f"EUR {val:,.2f}", f"EUR {hpl:+,.2f}"]
                y_row = 1.9 + ri * 0.55
                if ri % 2 == 0:
                    add_rect(s3, 0.3, y_row-0.05, 12.7, 0.5, (247,249,252))
                for ci, v in enumerate(row_vals):
                    col_color = (15,123,63) if ci == 4 and hpl >= 0 else (180,35,24) if ci == 4 else (45,58,74)
                    add_text(s3, v, x_starts[ci], y_row, col_w[ci], 0.45, 11, False, col_color)

        prs.save(str(path))
        return path
    except ImportError:
        logger.warning("python-pptx not installed, skipping PPTX generation")
        return None
    except Exception as e:
        logger.error("PPTX generation error: %s", e)
        return None


def _send_email(to: str, filepath: Optional[Path], fmt: str) -> bool:
    smtp_host = os.getenv("SMTP_HOST", "")
    smtp_user = os.getenv("SMTP_USER", "")
    smtp_pass = os.getenv("SMTP_PASS", "")
    from_email = os.getenv("FROM_EMAIL", smtp_user)
    smtp_port = int(os.getenv("SMTP_PORT", "587"))
    if not smtp_host or not smtp_user:
        logger.warning("SMTP not configured, skipping email")
        return False
    try:
        msg = MIMEMultipart()
        msg["From"] = from_email
        msg["To"] = to
        msg["Subject"] = "ETF Tracker — Report Portafoglio"
        msg.attach(MIMEText(
            "<h2>ETF Tracker</h2><p>In allegato trovi il report del tuo portafoglio.</p>"
            "<p><em>ETF Tracker non costituisce consulenza finanziaria.</em></p>",
            "html"
        ))
        if filepath and filepath.exists():
            with open(filepath, "rb") as f:
                part = MIMEBase("application", "octet-stream")
                part.set_payload(f.read())
            encoders.encode_base64(part)
            part.add_header("Content-Disposition", f'attachment; filename="{filepath.name}"')
            msg.attach(part)
        with smtplib.SMTP(smtp_host, smtp_port) as server:
            server.starttls()
            server.login(smtp_user, smtp_pass)
            server.send_message(msg)
        return True
    except Exception as e:
        logger.error("Email send error: %s", e)
        return False


# ── Static ETF data (inline — no external file dependency) ────────────────────

ETF_DATABASE = [
    {"isin":"IE00B4L5Y983","ticker":"VWCE","name":"Vanguard FTSE All-World","price":118.50,"chg1d":0.45,"chg1y":12.3,"ter":0.22,"aum":1850000,"asset":"Azionario","region":"Globale","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"Vanguard","rating":5,"chg5y":52.0},
    {"isin":"IE00B4L5YY86","ticker":"IWDA","name":"iShares Core MSCI World","price":85.20,"chg1d":0.32,"chg1y":14.1,"ter":0.20,"aum":980000,"asset":"Azionario","region":"Globale","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":5,"chg5y":55.0},
    {"isin":"IE00BDBRDM35","ticker":"EUNL","name":"iShares Core Euro STOXX 50","price":42.80,"chg1d":0.18,"chg1y":8.5,"ter":0.10,"aum":125000,"asset":"Azionario","region":"Europa","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":28.0},
    {"isin":"IE00BDBRDM43","ticker":"EXS1","name":"iShares STOXX Europe 600","price":95.40,"chg1d":0.25,"chg1y":9.2,"ter":0.20,"aum":85000,"asset":"Azionario","region":"Europa","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":30.0},
    {"isin":"IE00BKM4GZ66","ticker":"EMAE","name":"iShares Core MSCI EM IMI","price":32.15,"chg1d":0.85,"chg1y":6.8,"ter":0.18,"aum":145000,"asset":"Azionario","region":"Emergenti","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":18.0},
    {"isin":"IE00BDBRDM51","ticker":"IUSN","name":"iShares MSCI World Small Cap","price":8.95,"chg1d":0.42,"chg1y":11.5,"ter":0.35,"aum":42000,"asset":"Azionario","region":"Globale","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":35.0},
    {"isin":"IE00BDBRDM69","ticker":"IBGL","name":"iShares Euro Government Bond","price":136.20,"chg1d":-0.15,"chg1y":-2.1,"ter":0.09,"aum":185000,"asset":"Obbligazionario","region":"Europa","replication":"Fisica","distribution":"Distribuzione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":-8.0},
    {"isin":"IE00BDBRDM77","ticker":"IEGA","name":"iShares Core Euro Gov Bond","price":128.50,"chg1d":-0.12,"chg1y":-1.8,"ter":0.12,"aum":95000,"asset":"Obbligazionario","region":"Europa","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":-6.0},
    {"isin":"IE00BDBRDM85","ticker":"IBCI","name":"iShares Inflation Linked Govt","price":68.40,"chg1d":0.08,"chg1y":1.2,"ter":0.10,"aum":32000,"asset":"Obbligazionario","region":"Europa","replication":"Fisica","distribution":"Distribuzione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":5.0},
    {"isin":"IE00BDBRDM93","ticker":"XAD1","name":"Xtrackers MSCI AC World","price":98.75,"chg1d":0.38,"chg1y":13.5,"ter":0.19,"aum":65000,"asset":"Azionario","region":"Globale","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"DWS","rating":4,"chg5y":50.0},
    {"isin":"IE00BDBRDM01","ticker":"XBLC","name":"Xtrackers MSCI World Momentum","price":62.30,"chg1d":0.55,"chg1y":16.2,"ter":0.30,"aum":28000,"asset":"Azionario","region":"Globale","replication":"Campionamento","distribution":"Accumulazione","domicile":"Irlanda","issuer":"DWS","rating":4,"chg5y":58.0},
    {"isin":"IE00BDBRDM19","ticker":"XGLD","name":"Xetra-Gold","price":48.90,"chg1d":0.22,"chg1y":18.5,"ter":0.36,"aum":7800,"asset":"Commodities","region":"Globale","replication":"Fisica","distribution":"Accumulazione","domicile":"Germania","issuer":"Deutsche Borse","rating":3,"chg5y":42.0},
    {"isin":"IE00BDBRDM27","ticker":"IUSQ","name":"iShares MSCI ACWI","price":92.40,"chg1d":0.40,"chg1y":12.8,"ter":0.20,"aum":89000,"asset":"Azionario","region":"Globale","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":5,"chg5y":51.0},
]

MODEL_PORTFOLIOS = [
    {"id":"all-weather","name":"All-Weather","author":"Ray Dalio","risk":"Medio","riskLevel":3,"cagr":7.8,"maxDD":"-14.5","sharpe":0.72,"philosophy":"Asset allocation che performa in ogni condizione macroeconomica","allocation":[{"name":"Azionario","value":30,"color":"#0A2540"},{"name":"Obbligazionario","value":40,"color":"#1E5AA0"},{"name":"Oro","value":15,"color":"#D69E2E"},{"name":"Commodities","value":15,"color":"#68A063"}]},
    {"id":"bogleheads","name":"Bogleheads 3-Fund","author":"John Bogle","risk":"Medio","riskLevel":3,"cagr":8.2,"maxDD":"-18.3","sharpe":0.68,"philosophy":"Semplicita totale: azionario globale + obbligazionario + emergenti","allocation":[{"name":"Azionario Global","value":60,"color":"#0A2540"},{"name":"Emergenti","value":20,"color":"#1E5AA0"},{"name":"Obbligazionario","value":20,"color":"#8A94A6"}]},
    {"id":"permanent","name":"Permanent Portfolio","author":"Harry Browne","risk":"Basso","riskLevel":2,"cagr":5.4,"maxDD":"-8.1","sharpe":0.55,"philosophy":"Equa ripartizione per proteggere da ogni scenario economico","allocation":[{"name":"Azionario","value":25,"color":"#0A2540"},{"name":"Obbligazionario","value":25,"color":"#1E5AA0"},{"name":"Oro","value":25,"color":"#D69E2E"},{"name":"Cash","value":25,"color":"#8A94A6"}]},
    {"id":"coffeehouse","name":"Coffeehouse Portfolio","author":"Bill Schultheis","risk":"Basso","riskLevel":2,"cagr":6.8,"maxDD":"-12.4","sharpe":0.62,"philosophy":"Diversificazione estrema con 10 asset class","allocation":[{"name":"Large Cap","value":10,"color":"#0A2540"},{"name":"Small Value","value":10,"color":"#1E5AA0"},{"name":"REITs","value":10,"color":"#68A063"},{"name":"Gov Bond","value":10,"color":"#D69E2E"},{"name":"Emergenti","value":10,"color":"#8A94A6"},{"name":"Altro","value":50,"color":"#4A5568"}]},
    {"id":"golden","name":"Golden Butterfly","author":"Tyler","risk":"Medio-Basso","riskLevel":2,"cagr":7.1,"maxDD":"-11.2","sharpe":0.70,"philosophy":"Variante Permanent con tilt value e small cap","allocation":[{"name":"S&P 500","value":20,"color":"#0A2540"},{"name":"Small Cap Value","value":20,"color":"#1E3A5F"},{"name":"Long Gov Bond","value":20,"color":"#1E5AA0"},{"name":"Short Gov Bond","value":20,"color":"#8A94A6"},{"name":"Oro","value":20,"color":"#D69E2E"}]},
    {"id":"growth","name":"100% Growth","author":"ETF Tracker","risk":"Alto","riskLevel":4,"cagr":10.5,"maxDD":"-28.7","sharpe":0.58,"philosophy":"Massima esposizione azionaria per obiettivi di crescita","allocation":[{"name":"Azionario Global","value":70,"color":"#0A2540"},{"name":"Small Cap","value":20,"color":"#1E3A5F"},{"name":"Emergenti","value":10,"color":"#0F7B3F"}]},
]
